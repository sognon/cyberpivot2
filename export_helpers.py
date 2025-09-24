
# export_helpers.py
# Helpers to include Risks in DOCX reports and a "Top risques" slide in PPTX.

from typing import List
import pandas as pd

# -------- DOCX --------
def append_risks_section(doc, risks_df: pd.DataFrame) -> None:
    """
    Append a 'Risques' section into a python-docx Document.
    The section shows RiskID, Titre, Impact, Vraisemblance, Score, Classe, Contrôles liés.
    If 'Score' is missing, it will be computed as Impact * Vraisemblance.
    """
    doc.add_heading("Risques", level=1)
    if risks_df is None or risks_df.empty:
        doc.add_paragraph("Aucun risque défini.")
        return

    # prepare clean dataframe
    cols = ["RiskID","Titre","Impact","Vraisemblance","Score","Classe","Contrôles liés"]
    df = risks_df.copy()
    # Normalize numerics and compute Score if absent
    if "Score" not in df.columns:
        df["Impact"] = pd.to_numeric(df.get("Impact", 0), errors="coerce").fillna(0).astype(int)
        df["Vraisemblance"] = pd.to_numeric(df.get("Vraisemblance", 0), errors="coerce").fillna(0).astype(int)
        df["Score"] = (df["Impact"] * df["Vraisemblance"]).astype(int)

    # reindex columns
    for c in cols:
        if c not in df.columns:
            df[c] = ""

    # build table
    table = doc.add_table(rows=1, cols=len(cols))
    hdr = table.rows[0].cells
    for i, h in enumerate(cols):
        hdr[i].text = h

    for _, r in df[cols].iterrows():
        row = table.add_row().cells
        for i, c in enumerate(cols):
            row[i].text = str(r[c]) if pd.notna(r[c]) else ""

    # small spacing after
    doc.add_paragraph("")


# -------- PPTX --------
def add_top_risks_slide(prs, risks_df: pd.DataFrame, top_n: int = 5) -> None:
    """
    Add a slide with a 'Top risques' table to a python-pptx Presentation.
    The table shows RiskID, Titre, Score, Classe.
    """
    if risks_df is None or risks_df.empty:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Top risques"
        slide.placeholders[1].text = "Aucun risque défini."
        return

    df = risks_df.copy()
    # compute Score if needed
    if "Score" not in df.columns:
        df["Impact"] = pd.to_numeric(df.get("Impact", 0), errors="coerce").fillna(0).astype(int)
        df["Vraisemblance"] = pd.to_numeric(df.get("Vraisemblance", 0), errors="coerce").fillna(0).astype(int)
        df["Score"] = (df["Impact"] * df["Vraisemblance"]).astype(int)

    cols = ["RiskID","Titre","Score","Classe"]
    for c in cols:
        if c not in df.columns:
            df[c] = ""

    df = df.sort_values("Score", ascending=False).head(top_n)

    # Build slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    slide.shapes.title.text = "Top risques"

    # add table
    rows, cols_n = len(df) + 1, len(cols)
    left = Inches(0.5); top = Inches(1.8); width = Inches(9); height = Inches(0.8 + 0.3 * rows)

    table = slide.shapes.add_table(rows, cols_n, left, top, width, height).table
    # header
    for i, h in enumerate(cols):
        table.cell(0, i).text = h
    # rows
    for r_idx, (_, row) in enumerate(df.iterrows(), start=1):
        table.cell(r_idx, 0).text = str(row["RiskID"])
        table.cell(r_idx, 1).text = str(row["Titre"])
        table.cell(r_idx, 2).text = str(row["Score"])
        table.cell(r_idx, 3).text = str(row["Classe"])


# Convenience imports for pptx positioning
from pptx.util import Inches
