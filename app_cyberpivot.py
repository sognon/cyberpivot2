# -*- coding: utf-8 -*-
"""
CyberPivot — App Streamlit (ISACA+, persistance par projet, formulaire stable, DOCX soigné, PPTX premium)
Onglets: Infos client | Audit | Preuves | Synthèse | Exports | (Admin)
- Évaluation à 4 niveaux: Non applicable | Conforme | Partiellement conforme | Non conforme
- Formulaire d'évaluation sans rafraîchissements intempestifs (st.form + sauvegarde explicite)
- Rapport DOCX mieux rédigé et mis en forme (titres brandés, vraies puces & numéros, espacements)
- Présentation PowerPoint de restitution (si python-pptx installé)
- Admin: charger/éditer/activer des normes ; paramétrage des seuils & recos ; gestion minimale d’utilisateurs
- Persistance: data/projects/<audit_slug>/{meta.json,audit.csv,plan.csv}, preuves : data/evidence/<slug>/
"""

from __future__ import annotations
import io, re, json, zipfile, hashlib
from datetime import datetime, timedelta
from pathlib import Path
from tempfile import NamedTemporaryFile
from typing import Optional, Dict, Any, List

import numpy as np
import pandas as pd
import streamlit as st
import matplotlib.pyplot as plt
from PIL import Image

from docx import Document
from docx.shared import Pt, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn

# PowerPoint (optionnel)
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt as PPTPt
    from pptx.dml.color import RGBColor as PPTColor
except Exception:
    Presentation = None  # si non installé, on désactivera le bouton

# ---- Auth locale (gardez votre auth.py) ----
import auth

APP_BUILD = "ISACA_SV_v2.4-stableform+cleanDOCX+pptx"

# ---- Dossiers ----
DATA_DIR       = Path("./data")
NORMS_DIR      = DATA_DIR / "norms"
PROJECTS_DIR   = DATA_DIR / "projects"
EVIDENCE_DIR   = DATA_DIR / "evidence"
SETTINGS_PATH  = DATA_DIR / "settings.json"
ACTIVE_NORM    = DATA_DIR / "active_norm.txt"
RECENTS_PATH   = DATA_DIR / "recents.json"

for p in (DATA_DIR, NORMS_DIR, PROJECTS_DIR, EVIDENCE_DIR):
    p.mkdir(parents=True, exist_ok=True)

# ---- Colonnes & paramètres par défaut ----
REQUIRED_COLS = ["Domain", "ID", "Item", "Contrôle", "Level", "Comment"]

# Échelle d'évaluation — 4 choix (ordre figé, pas d'entrée vide)
EVAL_OPTIONS = ["Non applicable", "Conforme", "Partiellement conforme", "Non conforme"]

DEFAULT_SETTINGS = {
    "level_normalization": {
        "non applicable": "Non applicable", "na": "Non applicable", "n/a": "Non applicable",
        "conforme": "Conforme", "ok": "Conforme", "pass": "Conforme",
        "partiellement conforme": "Partiellement conforme", "partial": "Partiellement conforme", "partiel": "Partiellement conforme",
        "non conforme": "Non conforme", "ko": "Non conforme", "fail": "Non conforme",
        # anciens imports → on rabat sur 4 niveaux
        "faible": "Partiellement conforme", "low": "Partiellement conforme",
        "moyen": "Partiellement conforme", "medium": "Partiellement conforme",
        "élevé": "Non conforme", "eleve": "Non conforme", "high": "Non conforme",
        "critique": "Non conforme", "critical": "Non conforme",
        "non evalue": "Partiellement conforme", "non évalué": "Partiellement conforme", "ne": "Partiellement conforme",
        "": "Partiellement conforme", " ": "Partiellement conforme"
    },
    # échelle interne (pour chiffrage/risques)
    "risk_score": {"Élevé":75, "Moyen":50, "Conforme":0, "Non applicable":0},
    "deadlines_days": {"Élevé":60, "Moyen":90, "Conforme":180, "Non applicable":180},
    "effort_days": {"Élevé":10, "Moyen":5, "Conforme":1, "Non applicable":0},
    "budget_eur": {"Élevé":10000, "Moyen":5000, "Conforme":1000, "Non applicable":0},
    "resource_profiles": [
        {"pattern":"accès|identit|iam", "humans":"RSSI, Admin AD/IAM, Owners apps", "tech":"MFA, RBAC, PAM, IGA"},
        {"pattern":"opéra|ops|production", "humans":"Resp Ops/Prod, Admin Systèmes, SecOps", "tech":"EDR, SIEM, Patch mgmt, Backups immuables"},
        {"pattern":"fourn|tiers|supplier", "humans":"Achats, Juridique, RSSI, Contract manager", "tech":"Tier risk tooling, clauses sécurité"},
        {"pattern":"gouv|organis", "humans":"Direction, RSSI, DPO, Risk manager", "tech":"Politique SSI, référentiels, GRC"},
        {"pattern":"industri|ot|ics", "humans":"Resp OT, Ingénieur process, RSSI", "tech":"ISA/IEC 62443, bastion, inventaire OT"},
    ],
    "recs": {
        "default": {
            "gouv": "Renforcer la gouvernance SSI (rôles, responsabilités, comités, RACI).",
            "pil":  "Mettre en place un pilotage avec objectifs, feuille de route et revues périodiques.",
            "fou":  "Encadrer les fournisseurs (clauses sécurité, preuves, revues).",
            "arch": "Améliorer l’architecture: segmentation, durcissement, bastion, journaux centralisés.",
            "kpi":  "KPI: % actions dans les délais, % conformités, MTTD/MTTR, patch<30j, % comptes revus."
        },
        "patterns": [
            {"pattern":"accès|iam", "gouv":"Politique d’habilitation & cycle de vie identités.",
             "pil":"Recertifications trimestrielles, owners nommés.",
             "fou":"Accès prestataires JIT/JEA, clauses PAM.",
             "arch":"MFA généralisée, PAM, RBAC, séparation env.",
             "kpi":"% MFA, % applis recertifiées, comptes orphelins."},
            {"pattern":"opéra|ops", "gouv":"Runbooks ops/sécu, astreintes.",
             "pil":"KPIs patch/antivirus/EDR, revues hebdo événements.",
             "fou":"SLA sécu contractualisés.",
             "arch":"Backups 3-2-1 immuables, durcissement, EDR+SIEM.",
             "kpi":"patch<30j, couverture EDR, restauration OK."}
        ]
    },
    "brand_color_hex": "2F5597"
}

# ---- Clés d'état éditeur norme ----
NORM_EDIT_ALIAS_KEY  = "norm_editor_alias"
NORM_EDIT_DATA_KEY   = "norm_editor_df_data"
NORM_EDIT_WIDGET_KEY = "norm_editor_editor"

# ------------------------ Utilitaires ------------------------
def _slug(s: str) -> str:
    s = (s or "").strip().lower()
    s = re.sub(r"[^\w\s-]", "", s)
    s = re.sub(r"\s+", "-", s)
    return s[:80] or "item"

def _project_dir(slug: str) -> Path:
    p = PROJECTS_DIR / slug
    p.mkdir(parents=True, exist_ok=True)
    (p / "exports").mkdir(exist_ok=True)
    return p

def _project_paths(slug: str) -> Dict[str, Path]:
    base = _project_dir(slug)
    return {
        "dir": base,
        "meta": base / "meta.json",
        "audit_csv": base / "audit.csv",
        "plan_csv": base / "plan.csv",
        "exports_dir": base / "exports",
        "logo": base / "client_logo.png",
        "heatmap": base / "exports" / f"heatmap_{slug}.png",
        "radar":   base / "exports" / f"radar_{slug}.png",
        "docx":    base / "exports" / f"rapport_isaca_{slug}.docx",
        "pptx":    base / "exports" / f"restitution_{slug}.pptx",
    }

# ----- Projets récents -----
def _load_recents() -> Dict[str, str]:
    if RECENTS_PATH.exists():
        try: return json.loads(RECENTS_PATH.read_text(encoding="utf-8"))
        except Exception: return {}
    return {}

def _touch_recent(slug: str) -> None:
    r = _load_recents()
    r[_slug(slug)] = datetime.utcnow().isoformat()+"Z"
    r = dict(sorted(r.items(), key=lambda kv: kv[1], reverse=True)[:12])
    RECENTS_PATH.write_text(json.dumps(r, ensure_ascii=False, indent=2), encoding="utf-8")

def _recent_list(current: str) -> List[str]:
    r = _load_recents()
    order = [k for k, _ in sorted(r.items(), key=lambda kv: kv[1], reverse=True)]
    if current not in order:
        order = [current] + order
    return order

# ---- Helper toggle (initialisation sûre) ----
def _sidebar_toggle(label: str, key: str, default: bool = True) -> bool:
    if key not in st.session_state:
        st.session_state[key] = default
    return st.sidebar.toggle(label, key=key)

# ---- Signatures (autosave) ----
def _sig_csv(df: pd.DataFrame) -> str:
    return hashlib.sha1(df.to_csv(index=False).encode("utf-8")).hexdigest()

# ------------------------ Normalisation & risques ------------------------
def _norm_level(x: Any) -> str:
    if x is None: return "Partiellement conforme"
    t = str(x).strip()
    if t == "": return "Partiellement conforme"
    lo = t.lower()
    mapping = st.session_state.get("settings", DEFAULT_SETTINGS).get("level_normalization", DEFAULT_SETTINGS["level_normalization"])
    return mapping.get(lo, t)

def _ensure_df(df: pd.DataFrame) -> pd.DataFrame:
    d = df.copy()
    synonyms = {
        "Domain": ["domaine","domain"],
        "ID": ["identifiant","id","ref","reference"],
        "Item": ["item","intitulé","intitule","titre"],
        "Contrôle": ["controle","contrôle","control","objectif","exigence"],
        "Level": ["niveau","level","maturité","maturite","conformité","conformite","status"],
        "Comment": ["comment","commentaire","note","observation"]
    }
    for need in REQUIRED_COLS:
        if need not in d.columns:
            for syn in synonyms.get(need, []):
                for col in d.columns:
                    if col.strip().lower() == syn:
                        d[need] = d[col]; break
                if need in d.columns: break
            if need not in d.columns:
                d[need] = ""
        d[need] = d[need].astype(str)

    # force l'échelle 4 (pas d'option vide)
    d["Level"] = d["Level"].apply(lambda s: "Partiellement conforme" if str(s).strip()=="" else str(s))
    d["Level"] = d["Level"].map(_norm_level)
    d.loc[~d["Level"].isin(EVAL_OPTIONS), "Level"] = "Partiellement conforme"
    return d[REQUIRED_COLS]

def _risk_from_level(level: str) -> str:
    level = _norm_level(level)
    if level == "Non conforme": return "Élevé"
    if level == "Partiellement conforme": return "Moyen"
    return level  # Conforme / Non applicable

def _risk_score(level: str) -> int:
    scores = st.session_state.get("settings", DEFAULT_SETTINGS).get("risk_score", DEFAULT_SETTINGS["risk_score"])
    return int(scores.get(_risk_from_level(level), 0))

def _deadline_from_level(level: str) -> str:
    days_map = st.session_state.get("settings", DEFAULT_SETTINGS).get("deadlines_days", DEFAULT_SETTINGS["deadlines_days"])
    days = int(days_map.get(_risk_from_level(level), 90))
    return (datetime.utcnow() + timedelta(days=days)).date().isoformat()

# ------------------------ Settings I/O ------------------------
def load_settings() -> Dict[str, Any]:
    try:
        if SETTINGS_PATH.exists():
            data = json.loads(SETTINGS_PATH.read_text(encoding="utf-8"))
            base = json.loads(json.dumps(DEFAULT_SETTINGS))
            base.update(data)
            return base
    except Exception:
        pass
    return json.loads(json.dumps(DEFAULT_SETTINGS))

def save_settings(s: Dict[str, Any]) -> None:
    SETTINGS_PATH.write_text(json.dumps(s, ensure_ascii=False, indent=2), encoding="utf-8")

def export_settings_to_excel(settings: Dict[str, Any]) -> bytes:
    # Nécessite openpyxl
    buf = io.BytesIO()
    with pd.ExcelWriter(buf, engine="openpyxl") as w:
        pd.DataFrame(list(settings.get("level_normalization", {}).items()),
                     columns=["input","normalized"]).to_excel(w, "level_normalization", index=False)
        pd.DataFrame(list(settings.get("risk_score", {}).items()),
                     columns=["level","score"]).to_excel(w, "risk_score", index=False)
        pd.DataFrame(list(settings.get("deadlines_days", {}).items()),
                     columns=["level","days"]).to_excel(w, "deadlines_days", index=False)
        pd.DataFrame(list(settings.get("effort_days", {}).items()),
                     columns=["level","effort_jh"]).to_excel(w, "effort_days", index=False)
        pd.DataFrame(list(settings.get("budget_eur", {}).items()),
                     columns=["level","budget_eur"]).to_excel(w, "budget_eur", index=False)
        pd.DataFrame(settings.get("resource_profiles", [])).to_excel(w, "resource_profiles", index=False)
        recs = settings.get("recs", {})
        pd.DataFrame([recs.get("default", {})]).to_excel(w, "recs_default", index=False)
        pd.DataFrame(recs.get("patterns", [])).to_excel(w, "recs_patterns", index=False)
    return buf.getvalue()

def import_settings_from_excel(file) -> Dict[str, Any]:
    xls = pd.ExcelFile(file)
    s: Dict[str, Any] = json.loads(json.dumps(DEFAULT_SETTINGS))
    if "level_normalization" in xls.sheet_names:
        df = pd.read_excel(xls, "level_normalization")
        s["level_normalization"] = {str(r["input"]).strip().lower(): str(r["normalized"]).strip()
                                    for _, r in df.dropna(how="all").iterrows() if str(r.get("input","")).strip()}
    def sheet_to_dict(sheet, k1, k2, cast=int):
        if sheet in xls.sheet_names:
            df = pd.read_excel(xls, sheet)
            out={}
            for _, r in df.dropna(how="all").iterrows():
                a = str(r.get(k1,"")).strip(); b = r.get(k2, None)
                if a:
                    try: out[a] = cast(b) if b==b else None
                    except Exception: out[a] = b
            return {k:v for k,v in out.items() if v is not None}
        return {}
    s["risk_score"]     = sheet_to_dict("risk_score","level","score",int) or s["risk_score"]
    s["deadlines_days"] = sheet_to_dict("deadlines_days","level","days",int) or s["deadlines_days"]
    s["effort_days"]    = sheet_to_dict("effort_days","level","effort_jh",int) or s["effort_days"]
    s["budget_eur"]     = sheet_to_dict("budget_eur","level","budget_eur",int) or s["budget_eur"]
    if "resource_profiles" in xls.sheet_names:
        dfp = pd.read_excel(xls, "resource_profiles")
        s["resource_profiles"] = [dict(r.fillna("")) for _, r in dfp.iterrows()]
    recs = s.get("recs", {"default":{}, "patterns":[]})
    if "recs_default" in xls.sheet_names:
        rdf = pd.read_excel(xls, "recs_default")
        if not rdf.empty:
            recs["default"] = {k: str(v) for k, v in rdf.iloc[0].fillna("").items()}
    if "recs_patterns" in xls.sheet_names:
        rpf = pd.read_excel(xls, "recs_patterns")
        recs["patterns"] = [dict(r.fillna("")) for _, r in rpf.iterrows()]
    s["recs"] = recs
    return s

# ------------------------ Persistance projet ------------------------
def load_meta_for(slug: str) -> Dict[str, Any]:
    p = _project_paths(slug)["meta"]
    if p.exists():
        try: return json.loads(p.read_text(encoding="utf-8"))
        except Exception: pass
    return {"audit_id":"AUDIT-001","client":"Client X","auditor":"","date": datetime.utcnow().date().isoformat(),"norm":_get_active_norm()}

def save_meta_for(slug: str, meta: Dict[str, Any]) -> None:
    paths = _project_paths(slug)
    paths["meta"].write_text(json.dumps(meta, ensure_ascii=False, indent=2), encoding="utf-8")
    _touch_recent(slug)

def load_df_for(slug: str, meta: Optional[Dict[str,Any]]=None) -> pd.DataFrame:
    csvp = _project_paths(slug)["audit_csv"]
    if csvp.exists():
        try: return _ensure_df(pd.read_csv(csvp))
        except Exception: pass
    norm = (meta or {}).get("norm") or _get_active_norm()
    try:
        return _ensure_df(_load_norm_df(norm)) if norm else _load_active_df()
    except Exception:
        return _load_active_df()

def save_df_for(slug: str, df: pd.DataFrame) -> None:
    _ensure_df(df).to_csv(_project_paths(slug)["audit_csv"], index=False)
    _touch_recent(slug)

def load_plan_for(slug: str, df: Optional[pd.DataFrame]=None) -> pd.DataFrame:
    p = _project_paths(slug)["plan_csv"]
    if p.exists():
        try: return pd.read_csv(p)
        except Exception: pass
    d = df if df is not None else load_df_for(slug)
    plan = build_action_plan(d); save_plan_for(slug, plan); return plan

def save_plan_for(slug: str, plan: pd.DataFrame) -> None:
    plan.to_csv(_project_paths(slug)["plan_csv"], index=False)
    _touch_recent(slug)

def list_projects() -> List[str]:
    return sorted([d.name for d in PROJECTS_DIR.iterdir() if d.is_dir()])

# ------------------------ Normes ------------------------
def _save_norm_df(name: str, df: pd.DataFrame) -> Path:
    name = _slug(name)
    p = NORMS_DIR / f"{name}.csv"
    _ensure_df(df).to_csv(p, index=False)
    return p

def _load_norm_df(name: str) -> pd.DataFrame:
    p_csv = NORMS_DIR / f"{_slug(name)}.csv"
    if p_csv.exists():
        return _ensure_df(pd.read_csv(p_csv))
    raise FileNotFoundError(f"Norme '{name}' introuvable.")

def _list_norms() -> List[str]:
    return sorted([p.stem for p in NORMS_DIR.glob("*.csv")])

def _set_active_norm(name: str) -> None:
    ACTIVE_NORM.write_text(_slug(name), encoding="utf-8")

def _get_active_norm() -> Optional[str]:
    if ACTIVE_NORM.exists():
        return ACTIVE_NORM.read_text(encoding="utf-8").strip()
    return None

def _load_active_df() -> pd.DataFrame:
    active = _get_active_norm()
    if active:
        try: return _load_norm_df(active)
        except Exception: pass
    # échantillon de départ
    return _ensure_df(pd.DataFrame([
        {"Domain":"Gouvernance","ID":"GOV-01","Item":"Charte sécurité","Contrôle":"Existence charte","Level":"Non conforme","Comment":""},
        {"Domain":"Opérations","ID":"OPS-02","Item":"Sauvegardes","Contrôle":"Plan de sauvegarde","Level":"Partiellement conforme","Comment":"Pas de PRA testé"},
        {"Domain":"Accès","ID":"ACC-03","Item":"Gestion des comptes","Contrôle":"Revue périodique des accès","Level":"Non conforme","Comment":"Absence de revue formelle"},
    ]))

# ------------------------ Plan d’actions ------------------------
def _default_resources(row: pd.Series) -> tuple[str, str, int, int]:
    settings = st.session_state.get("settings", DEFAULT_SETTINGS)
    dom = str(row.get("Domain","")).lower()
    risk = _risk_from_level(row.get("Level",""))
    effort = int(settings.get("effort_days", DEFAULT_SETTINGS["effort_days"]).get(risk, 3))
    budget = int(settings.get("budget_eur",  DEFAULT_SETTINGS["budget_eur"]).get(risk, 3000))
    humans = "RSSI, Équipe IT/Sécu, Métier concerné"
    tech   = "Outillage (EDR/SIEM/MFA/Backups)"
    for prof in settings.get("resource_profiles", DEFAULT_SETTINGS["resource_profiles"]):
        try:
            if re.search(prof.get("pattern",""), dom, flags=re.IGNORECASE):
                humans = prof.get("humans", humans); tech = prof.get("tech", tech); break
        except Exception:
            continue
    return humans, tech, effort, budget

def _default_recs(row: pd.Series) -> tuple[str, str, str, str, str]:
    settings = st.session_state.get("settings", DEFAULT_SETTINGS)
    dom = str(row.get("Domain","")).lower()
    rec = settings.get("recs", DEFAULT_SETTINGS["recs"])
    base = rec.get("default", {})
    gouv = base.get("gouv",""); pil = base.get("pil",""); fou = base.get("fou",""); arch = base.get("arch",""); kpi = base.get("kpi","")
    for pat in rec.get("patterns", []):
        try:
            if re.search(pat.get("pattern",""), dom, flags=re.IGNORECASE):
                gouv = pat.get("gouv", gouv); pil = pat.get("pil", pil); fou  = pat.get("fou", fou)
                arch = pat.get("arch", arch); kpi = pat.get("kpi", kpi); break
        except Exception:
            continue
    return gouv, pil, fou, arch, kpi

def build_action_plan(df: pd.DataFrame) -> pd.DataFrame:
    d = _ensure_df(df)
    d["RiskLevel"] = d["Level"].map(_risk_from_level)
    d["RiskScore"] = d["RiskLevel"].map(_risk_score)
    sel = d[d["RiskScore"] > 0].copy()
    if sel.empty:
        return pd.DataFrame(columns=[
            "Domaine","ID","Intitulé","Contrôle","Niveau de risque","Score",
            "Action recommandée","Responsable","Échéance","Statut",
            "Effort (JH)","Ressources humaines","Ressources techniques","Budget estimé (€)",
            "Reco — Gouvernance","Reco — Pilotage","Reco — Fournisseurs","Reco — Architecture",
            "Avancement (%)","Indicateurs de suivi"
        ])
    rows=[]
    for _, r in sel.iterrows():
        rh, tech, effort, budget = _default_resources(r)
        rec_g, rec_p, rec_f, rec_a, kpi = _default_recs(r)
        rows.append({
            "Domaine": r["Domain"], "ID": r["ID"], "Intitulé": r["Item"], "Contrôle": r["Contrôle"],
            "Niveau de risque": r["RiskLevel"], "Score": r["RiskScore"],
            "Action recommandée": f"Mettre en conformité : {(r.get('Contrôle') or r.get('Item') or 'contrôle')}",
            "Responsable": "", "Échéance": _deadline_from_level(r["Level"]), "Statut": "À faire",
            "Effort (JH)": effort, "Ressources humaines": rh, "Ressources techniques": tech, "Budget estimé (€)": budget,
            "Reco — Gouvernance": rec_g, "Reco — Pilotage": rec_p, "Reco — Fournisseurs": rec_f, "Reco — Architecture": rec_a,
            "Avancement (%)": 0, "Indicateurs de suivi": kpi,
        })
    return pd.DataFrame(rows)

# ------------------------ DOCX helpers/report ------------------------
def _set_cell_bg(cell, hex_color: str):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    shd = OxmlElement('w:shd')
    shd.set(qn('w:fill'), hex_color)
    tcPr.append(shd)

def _brand_hex():
    return st.session_state.get("settings", DEFAULT_SETTINGS).get("brand_color_hex", "2F5597")

def _apply_doc_branding(doc: Document, meta: Dict[str, Any]):
    brand_hex = _brand_hex()
    for section in doc.sections:
        header = section.header
        p = header.paragraphs[0] if header.paragraphs else header.add_paragraph()
        p.text = f"{meta.get('client','')} — {meta.get('audit_id','')}"
        if p.runs:
            p.runs[0].font.color.rgb = RGBColor.from_string(brand_hex)

def _style_normal(doc: Document):
    try:
        n = doc.styles["Normal"]
        n.font.name = "Calibri"
        n.font.size = Pt(11)
    except Exception:
        pass

def _heading(doc: Document, text: str, level: int = 1):
    p = doc.add_heading(text, level=level)
    try:
        for r in p.runs:
            r.font.color.rgb = RGBColor.from_string(_brand_hex())
    except Exception:
        pass
    return p

def _para(doc: Document, text: str, bold: bool = False, space_after_pt: int = 6):
    p = doc.add_paragraph(text)
    if bold and p.runs:
        p.runs[0].bold = True
    p.paragraph_format.space_after = Pt(space_after_pt)
    return p

def _bullets(doc: Document, items: list[str], indent_cm: float = 0.5, space_after_pt: int = 2):
    for t in items:
        p = doc.add_paragraph(t, style="List Bullet")
        p.paragraph_format.left_indent = Cm(indent_cm)
        p.paragraph_format.space_after = Pt(space_after_pt)

def _numbers(doc: Document, items: list[str], indent_cm: float = 0.5, space_after_pt: int = 2):
    for t in items:
        p = doc.add_paragraph(t, style="List Number")
        p.paragraph_format.left_indent = Cm(indent_cm)
        p.paragraph_format.space_after = Pt(space_after_pt)

def _counts_text_4lvl(d: pd.DataFrame) -> str:
    counts = d["Level"].value_counts().to_dict(); C = lambda n:int(counts.get(n,0))
    return (f"Non applicable : {C('Non applicable')} • Conforme : {C('Conforme')} • "
            f"Partiellement conforme : {C('Partiellement conforme')} • Non conforme : {C('Non conforme')}")

def generate_isaca_docx(audit_id: str, df: pd.DataFrame, meta: Dict[str, Any], plan: pd.DataFrame,
                        figs: Dict[str, Path], out_path: Path, logo_path: Optional[Path] = None) -> Path:
    d = _ensure_df(df)
    doc = Document()
    _style_normal(doc)

    # Marges
    for s in doc.sections:
        s.top_margin = Cm(2); s.bottom_margin = Cm(2); s.left_margin = Cm(2); s.right_margin = Cm(2)

    # Propriétés
    doc.core_properties.title = f"Rapport d'audit — {audit_id}"
    doc.core_properties.author = meta.get("auditor","")
    doc.core_properties.company = meta.get("client","")

    # Page de garde
    tbl = doc.add_table(rows=1, cols=2)
    left, right = tbl.rows[0].cells
    p = left.paragraphs[0]; p.alignment = WD_ALIGN_PARAGRAPH.LEFT
    r = p.add_run("Rapport d'audit — ISACA\n"); r.bold = True; r.font.size = Pt(24); r.font.color.rgb = RGBColor.from_string(_brand_hex())
    p.add_run(f"{audit_id}\n").font.size = Pt(14)
    p.add_run(f"Client : {meta.get('client','')}\n")
    p.add_run(f"Date : {meta.get('date', datetime.utcnow().date().isoformat())}\n")
    p.add_run(f"Auditeur : {meta.get('auditor','')}\n")
    if logo_path and Path(logo_path).exists():
        rp = right.paragraphs[0]; rp.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        rp.add_run().add_picture(str(logo_path), width=Cm(4))
    doc.add_page_break()

    _apply_doc_branding(doc, meta)

    # 1. Résumé exécutif
    _heading(doc, "1. Résumé exécutif", 1)
    _para(doc, "Synthèse globale.", bold=True, space_after_pt=2)
    _para(doc,
          f"L’évaluation a porté sur {len(d)} contrôles. Répartition des statuts : {_counts_text_4lvl(d)}.",
          space_after_pt=6)

    # principaux domaines à risque
    def _top_domains_text(dd: pd.DataFrame, n=3) -> str:
        if dd.empty: return "n/a"
        tmp = dd.copy()
        tmp["RiskScore"] = tmp["Level"].map(_risk_from_level).map(_risk_score)
        by = tmp.groupby("Domain")["RiskScore"].sum().sort_values(ascending=False).head(n)
        return ", ".join([f"{k} (score {int(v)})" for k, v in by.items()]) or "n/a"
    _para(doc, f"Les écarts les plus significatifs concernent : {_top_domains_text(d, 3)}.", space_after_pt=6)

    _para(doc, "Priorités de remédiation :", bold=True, space_after_pt=2)
    _numbers(doc, [
        "Traiter en priorité les non-conformités à fort impact et forte exposition.",
        "Planifier la mise en conformité des points partiellement conformes avec jalons réalistes.",
        "Sécuriser les dépendances fournisseurs et les accès à privilèges."
    ], indent_cm=0.8, space_after_pt=6)

    # 2. Méthodologie
    _heading(doc, "2. Méthodologie", 1)
    _bullets(doc, [
        "Périmètre et référentiel : grille d’audit basée sur la norme sélectionnée et adaptée au contexte.",
        "Normalisation : toute valeur importée est ramenée à l’échelle ci-dessous pour permettre comparaison et consolidation."
    ], indent_cm=0.5)

    _para(doc, "Échelle d’évaluation (4 niveaux) :", space_after_pt=2)
    _bullets(doc, [
        "Non applicable — le contrôle ne s’applique pas au périmètre étudié.",
        "Conforme — exigences respectées et preuves suffisantes.",
        "Partiellement conforme — exigences partiellement atteintes ; écart à résorber.",
        "Non conforme — exigences non respectées ; écart majeur."
    ], indent_cm=1.0, space_after_pt=4)

    _bullets(doc, [
        "Cartographie des risques : Non conforme → « Élevé » ; Partiellement conforme → « Moyen » ; "
        "Conforme et Non applicable → 0. Les scores sont agrégés par domaine.",
        "Plan d’actions : pour chaque écart, une action avec responsable, échéance, effort (JH) et budget. "
        "Des profils de ressources types et des recommandations (gouvernance, pilotage, fournisseurs, architecture) accompagnent la mise en œuvre.",
        "Preuves : éléments collectés archivés par contrôle, exportables en ZIP avec manifeste.",
        "Limites : résultats basés sur les preuves disponibles et les entretiens à la date de l’audit."
    ], indent_cm=0.5, space_after_pt=6)

    # 3. Graphiques
    _heading(doc, "3. Cartographie et radar", 1)
    if figs.get("heatmap") and Path(figs["heatmap"]).exists():
        _para(doc, "Cartographie des risques (somme des scores par domaine)", space_after_pt=2)
        doc.add_picture(str(figs["heatmap"]), width=Cm(16))
    if figs.get("radar") and Path(figs["radar"]).exists():
        _para(doc, "Radar par domaine (plus haut = meilleur)", space_after_pt=2)
        doc.add_picture(str(figs["radar"]), width=Cm(16))

    # 4. Détails
    _heading(doc, "4. Résultats détaillés", 1)
    t = doc.add_table(rows=1, cols=6)
    t.style = "Light Grid Accent 1" if "Light Grid Accent 1" in [s.name for s in doc.styles] else "Table Grid"
    headers = ["Domaine","ID","Intitulé","Contrôle/objectif","Niveau","Constatations"]
    for i, col in enumerate(headers): t.rows[0].cells[i].text = col
    risk_color = {"Élevé":"FFC000","Moyen":"FFD966","Conforme":"00B050","Non applicable":"D9D9D9"}
    for _, row in d.iterrows():
        c = t.add_row().cells
        c[0].text = str(row["Domain"]); c[1].text = str(row["ID"]); c[2].text = str(row["Item"])
        c[3].text = str(row["Contrôle"]); c[4].text = str(row["Level"]); c[5].text = str(row["Comment"])
        _set_cell_bg(c[4], risk_color.get(_risk_from_level(row["Level"]), "FFFFFF"))

    # 5-6. Plan d'actions et ressources
    _heading(doc, "5. Plan d'actions priorisé", 1)
    if plan is not None and not plan.empty:
        cols_main = ["Domaine","ID","Intitulé","Niveau de risque","Score",
                     "Action recommandée","Responsable","Échéance","Statut",
                     "Effort (JH)","Budget estimé (€)","Avancement (%)"]
        t2 = doc.add_table(rows=1, cols=len(cols_main))
        t2.style = "Light List Accent 1" if "Light List Accent 1" in [s.name for s in doc.styles] else "Table Grid"
        for i, col in enumerate(cols_main): t2.rows[0].cells[i].text = col
        plan_sorted = plan.sort_values(by=["Niveau de risque","Score"], ascending=[True, False])
        for _, row in plan_sorted.iterrows():
            c = t2.add_row().cells
            for i, col in enumerate(cols_main): c[i].text = str(row.get(col,""))
            _set_cell_bg(c[cols_main.index("Niveau de risque")],
                         risk_color.get(str(row.get("Niveau de risque","")), "FFFFFF"))

        _heading(doc, "6. Ressources et recommandations", 1)
        cols_more = ["Domaine","ID","Ressources humaines","Ressources techniques",
                     "Reco — Gouvernance","Reco — Pilotage","Reco — Fournisseurs","Reco — Architecture",
                     "Indicateurs de suivi"]
        t3 = doc.add_table(rows=1, cols=len(cols_more))
        t3.style = "Light Grid Accent 2" if "Light Grid Accent 2" in [s.name for s in doc.styles] else "Table Grid"
        for i, col in enumerate(cols_more): t3.rows[0].cells[i].text = col
        for _, row in plan_sorted.iterrows():
            c = t3.add_row().cells
            for i, col in enumerate(cols_more): c[i].text = str(row.get(col,""))
    else:
        _para(doc, "Aucune action requise.", space_after_pt=6)

    doc.save(out_path); return out_path

# ------------------------ PowerPoint ------------------------
def generate_pptx_summary(audit_id: str, meta: Dict[str, Any], df: pd.DataFrame, plan: pd.DataFrame,
                          figs: Dict[str, Path], out_path: Path, brand_hex: str = "2F5597") -> Optional[Path]:
    if Presentation is None:
        return None
    prs = Presentation()
    brand_rgb = tuple(int(brand_hex[i:i+2],16) for i in (0,2,4))

    # Titre
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.add_textbox(Inches(0.7), Inches(0.6), Inches(8.2), Inches(1.2)).text_frame
    run = title.paragraphs[0].add_run()
    run.text = f"Restitution — Audit {audit_id}"
    run.font.size = PPTPt(44); run.font.bold = True; run.font.color.rgb = PPTColor(*brand_rgb)

    tf = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(8.2), Inches(2.0)).text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = f"Client : {meta.get('client','')}"
    for line in [f"Date : {meta.get('date', datetime.utcnow().date().isoformat())}",
                 f"Auditeur : {meta.get('auditor','')}",
                 f"Référentiel : {meta.get('norm','') or '(n/a)'}"]:
        p = tf.add_paragraph(); p.text = line

    # Résumé
    counts = df["Level"].value_counts()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    left = Inches(0.7); top = Inches(0.6)
    box = slide.shapes.add_textbox(left, top, Inches(8.2), Inches(1.0)).text_frame
    box.paragraphs[0].text = "Résumé exécutif"
    box.paragraphs[0].runs[0].font.size = PPTPt(28); box.paragraphs[0].runs[0].font.bold = True
    stats = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8.2), Inches(2.0)).text_frame
    stats.word_wrap = True
    lines = [
        f"Contrôles évalués : {len(df)}",
        f"Répartition : N/A {int(counts.get('Non applicable',0))} • Conforme {int(counts.get('Conforme',0))} • "
        f"Partiel {int(counts.get('Partiellement conforme',0))} • Non conforme {int(counts.get('Non conforme',0))}",
        "Priorités : 1) non-conformités majeures, 2) points partiels à fort impact, 3) tierces parties et comptes à privilèges."
    ]
    for i, line in enumerate(lines):
        (stats.paragraphs[i] if i < len(stats.paragraphs) else stats.add_paragraph()).text = line

    # Heatmap + Radar
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(Inches(0.7), Inches(0.6), Inches(8.2), Inches(0.6)).text_frame.text = "Cartographie"
    if figs.get("heatmap") and Path(figs["heatmap"]).exists():
        slide.shapes.add_picture(str(figs["heatmap"]), Inches(0.7), Inches(1.2), width=Inches(4.0))
    if figs.get("radar") and Path(figs["radar"]).exists():
        slide.shapes.add_picture(str(figs["radar"]), Inches(5.0), Inches(1.2), width=Inches(4.0))

    # Top actions
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(Inches(0.7), Inches(0.6), Inches(8.2), Inches(0.6)).text_frame.text = "Top actions"
    tbl_rows = min(8, len(plan))
    if tbl_rows > 0:
        table = slide.shapes.add_table(tbl_rows+1, 5, Inches(0.7), Inches(1.2), Inches(8.2), Inches(4.0)).table
        headers = ["ID","Intitulé","Niveau","Échéance","Responsable"]
        for j,h in enumerate(headers): table.cell(0,j).text = h
        plan_sorted = plan.sort_values(by=["Niveau de risque","Score"], ascending=[True, False]).head(tbl_rows)
        for i, (_, r) in enumerate(plan_sorted.iterrows(), start=1):
            table.cell(i,0).text = str(r.get("ID",""))
            table.cell(i,1).text = str(r.get("Intitulé",""))
            table.cell(i,2).text = str(r.get("Niveau de risque",""))
            table.cell(i,3).text = str(r.get("Échéance",""))
            table.cell(i,4).text = str(r.get("Responsable",""))

    prs.save(out_path)
    return out_path

# ------------------------ Auth adapter ------------------------
class AuthAdapter:
    @staticmethod
    def init_db() -> None:
        try: auth.init_auth_db()
        except Exception: pass

    @staticmethod
    def verify(email: str, pwd: str) -> bool:
        try: return bool(auth.verify_password(email, pwd))
        except Exception: return False

    @staticmethod
    def get_user(email: str) -> Optional[Dict[str, Any]]:
        try:
            u = auth.get_user(email)
            return u if isinstance(u, dict) else None
        except Exception:
            return None

    @staticmethod
    def get_or_create(email: str, full_name: Optional[str] = None, role: str = "user") -> Optional[Dict[str, Any]]:
        try:
            u = auth.get_or_create_user(email=email, full_name=full_name or email, role=role)
            if isinstance(u, dict) and u.get("email"):
                return u
            return AuthAdapter.get_user(email)
        except Exception:
            return None

    @staticmethod
    def create_user(email: str, password: str, full_name: str, role: str = "user", is_active: bool = True):
        try:
            return auth.create_user(email=email, password=password, full_name=full_name, role=role, is_active=is_active)
        except Exception as e:
            return (False, f"Erreur création: {e}")

    @staticmethod
    def login_get_user(email: str, pwd: str) -> Optional[Dict[str, Any]]:
        if not AuthAdapter.verify(email, pwd): return None
        u = AuthAdapter.get_user(email)
        if not (isinstance(u, dict) and u.get("email")):
            u = AuthAdapter.get_or_create(email=email, full_name=email, role="user")
        return u if (isinstance(u, dict) and u.get("email")) else None

# ------------------------ UI helpers ------------------------
def _topbar(user: Dict[str, Any], slug: str):
    recents = _recent_list(slug)
    cols = st.columns([5,3,2,1])
    cols[0].markdown(f"**Connecté :** {user.get('email')} — {user.get('full_name')}  _(rôle: {user.get('role','user')})_")
    new_sel = cols[1].selectbox("Projet", recents, index=recents.index(slug), key="top_proj_select")
    if new_sel != slug:
        st.session_state["project_slug"] = new_sel
        st.rerun()
    if cols[2].button("Dupliquer projet", key="btn_dup_proj"):
        new_slug = f"{slug}-copy"
        paths_old = _project_paths(slug); paths_new = _project_paths(new_slug)
        if paths_old["meta"].exists(): paths_new["meta"].write_bytes(paths_old["meta"].read_bytes())
        if paths_old["audit_csv"].exists(): paths_new["audit_csv"].write_bytes(paths_old["audit_csv"].read_bytes())
        if paths_old["plan_csv"].exists(): paths_new["plan_csv"].write_bytes(paths_old["plan_csv"].read_bytes())
        if paths_old["logo"].exists(): paths_new["logo"].write_bytes(paths_old["logo"].read_bytes())
        _touch_recent(new_slug)
        st.session_state["project_slug"] = new_slug
        st.success(f"Projet dupliqué → {new_slug}"); st.rerun()
    if cols[3].button("Se déconnecter", key="logout_btn"):
        st.session_state.pop("_auth", None); st.rerun()

# ------------------------ Onglets ------------------------
def _tab_infos_client(slug: str):
    st.subheader("🧾 Infos client")
    paths = _project_paths(slug)
    meta = load_meta_for(slug)

    c1, c2, c3 = st.columns(3)
    new_audit = c1.text_input("Identifiant d'audit", value=meta.get("audit_id","AUDIT-001"), key=f"info_audit_id_{slug}")
    meta["client"]   = c2.text_input("Client",   value=meta.get("client","Client X"), key=f"info_client_{slug}")
    meta["auditor"]  = c3.text_input("Auditeur", value=meta.get("auditor",""), key=f"info_auditor_{slug}")
    meta["date"]     = st.date_input("Date du rapport",
                                     value=pd.to_datetime(meta.get("date", datetime.utcnow().date())).date(),
                                     key=f"info_date_{slug}").isoformat()

    norms = _list_norms()
    default_norm = meta.get("norm") if meta.get("norm") in norms else (_get_active_norm() if _get_active_norm() in norms else (norms[0] if norms else None))
    meta["norm"] = st.selectbox("Norme de travail", norms if norms else ["(aucune)"],
                                index=(norms.index(default_norm) if (norms and default_norm in norms) else 0),
                                key=f"info_norm_select_{slug}") if norms else None

    logo_file = st.file_uploader("Logo client (PNG/JPG)", type=["png","jpg","jpeg"], key=f"info_logo_{slug}")
    cols = st.columns([1.5,1.5,2,2])
    if cols[0].button("💾 Enregistrer (projet)", key=f"info_save_btn_{slug}"):
        new_slug = _slug(new_audit)
        meta["audit_id"] = new_audit
        save_meta_for(new_slug, meta)
        if "df" in st.session_state: save_df_for(new_slug, st.session_state["df"])
        if "plan_actions" in st.session_state: save_plan_for(new_slug, st.session_state["plan_actions"])
        st.session_state["project_slug"] = new_slug
        st.success(f"Projet sauvegardé: {new_slug}")
        st.rerun()

    if cols[1].button("📚 Appliquer la norme", key=f"info_apply_norm_btn_{slug}"):
        try:
            if meta.get("norm"):
                df = _load_norm_df(meta["norm"])
                st.session_state["df"] = _ensure_df(df)
                save_df_for(slug, st.session_state["df"])
                st.success("Norme appliquée et enregistrée."); st.rerun()
            else:
                st.warning("Aucune norme sélectionnée.")
        except Exception as e:
            st.error(f"Impossible d'appliquer la norme : {e}")

    if cols[2].button("🖼️ Mettre à jour le logo", key=f"info_logo_btn_{slug}"):
        if logo_file:
            image = Image.open(logo_file); image.save(paths["logo"])
            st.success("Logo client enregistré.")
        else:
            st.warning("Sélectionnez un fichier image.")

    if cols[3].button("🔁 Recharger les infos du projet", key=f"info_reload_btn_{slug}"):
        st.rerun()

    if paths["logo"].exists():
        st.image(str(paths["logo"]), caption=f"Logo — {meta.get('client','')}", use_column_width=False, width=220)

    st.caption(f"Norme projet: **{meta.get('norm','(non définie)')}** • Dataset: `{paths['audit_csv'].name}`")

def _tab_audit(slug: str):
    st.subheader("✏️ Audit — Réponses aux contrôles (formulaire)")
    df = _ensure_df(st.session_state.get("df", load_df_for(slug)))
    alias = load_meta_for(slug).get("norm") or _get_active_norm() or "(inconnu)"
    st.caption(f"Norme en cours : **{alias}** — {len(df)} contrôles")

    # Filtres persistants
    f1, f2, f3 = st.columns(3)
    domain = f1.selectbox("Domaine", ["(Tous)"] + sorted(df["Domain"].unique().tolist()), key=f"audit_domain_{slug}")
    search = f2.text_input("Recherche (texte)", key=f"audit_search_{slug}")
    todo   = f3.checkbox("Afficher uniquement les écarts", key=f"audit_todo_{slug}")

    view = df.copy()
    if domain != "(Tous)": view = view[view["Domain"] == domain]
    if search:
        s = search.lower()
        view = view[view.apply(lambda r: s in " ".join(map(str, r.values)).lower(), axis=1)]
    if todo:
        view["RiskScore"] = view["Level"].map(_risk_score)
        view = view[view["RiskScore"] > 0]

    # Action de masse
    mass_cols = st.columns([2,1,1])
    set_to = mass_cols[0].selectbox("Régler le niveau (sur le filtre affiché)", ["(aucun)"] + EVAL_OPTIONS, key=f"audit_mass_set_{slug}")
    if mass_cols[1].button("Appliquer au filtre", key=f"audit_mass_apply_{slug}"):
        if set_to != "(aucun)":
            df.loc[view.index, "Level"] = set_to
            st.session_state["df"] = _ensure_df(df)
            save_df_for(slug, st.session_state["df"])
            st.success(f"Niveau réglé sur {len(view)} lignes ({set_to}).")
            st.rerun()

    # -------- Formulaire stable (sauvegarde explicite) --------
    with st.form(key=f"audit_form_{slug}", clear_on_submit=False):
        edited = st.data_editor(
            view,
            column_config={
                "Level": st.column_config.SelectboxColumn(options=EVAL_OPTIONS, required=True),
                "Comment": st.column_config.TextColumn(),
            },
            use_container_width=True,
            num_rows="dynamic",
            key=f"audit_editor_df_{slug}"
        )
        submitted = st.form_submit_button("💾 Enregistrer les modifications", type="primary")
        if submitted:
            base = _ensure_df(st.session_state.get("df", df)).set_index(["Domain","ID","Item"]).copy()
            edited_idx = _ensure_df(edited).set_index(["Domain","ID","Item"]).copy()
            common = base.index.intersection(edited_idx.index)
            if len(common) > 0:
                cols_upd = ["Level","Comment","Contrôle","Domain","ID","Item"]
                inter = [c for c in cols_upd if c in base.columns and c in edited_idx.columns]
                base.loc[common, inter] = edited_idx.loc[common, inter]
            new_rows = edited_idx.loc[~edited_idx.index.isin(base.index)]
            if not new_rows.empty:
                add = new_rows.reindex(columns=base.columns, fill_value="")
                base = pd.concat([base, add], axis=0)
            base = _ensure_df(base.reset_index()).drop_duplicates(subset=["Domain","ID","Item"], keep="last")
            st.session_state["df"] = base
            save_df_for(slug, base)
            st.success("Modifications enregistrées ✅")

def _evi_dir(audit_id: str, qid: str, item: str) -> Path:
    return EVIDENCE_DIR / _slug(audit_id) / f"{_slug(qid)}_{_slug(item)}"

def persist_uploads(audit_id: str, qid: str, item: str, files: List[Any]) -> int:
    d = _evi_dir(audit_id, qid, item); d.mkdir(parents=True, exist_ok=True)
    n=0
    for f in files or []:
        p = d / f.name
        with open(p,"wb") as out: out.write(f.read())
        n+=1
    return n

def list_evidence(audit_id: str, qid: str, item: str) -> List[Path]:
    d = _evi_dir(audit_id, qid, item)
    return [p for p in d.iterdir() if p.is_file()] if d.exists() else []

def delete_evidence(path: Path) -> None:
    try: path.unlink(missing_ok=True)
    except Exception: pass

def export_evidence_zip(audit_id: str, df: pd.DataFrame) -> bytes:
    slug_audit = _slug(audit_id)
    with NamedTemporaryFile(delete=False, suffix=f"_{slug_audit}.zip") as tmp:
        with zipfile.ZipFile(tmp.name, "w", zipfile.ZIP_DEFLATED) as z:
            manifest=[]
            for _, row in df.iterrows():
                qid, item = row.get("ID",""), row.get("Item","")
                for p in list_evidence(audit_id, qid, item):
                    arc = f"{slug_audit}/{_slug(qid)}_{_slug(item)}/{p.name}"
                    z.write(p, arcname=arc); manifest.append({"id":qid,"item":item,"file":arc})
            csv_buf = io.StringIO(); pd.DataFrame(manifest).to_csv(csv_buf, index=False)
            z.writestr(f"{slug_audit}/manifest.csv", csv_buf.getvalue())
            z.writestr(f"{slug_audit}/manifest.json", json.dumps(manifest, ensure_ascii=False, indent=2))
        return Path(tmp.name).read_bytes()

def _tab_preuves(slug: str):
    st.subheader("📎 Preuves")
    df = _ensure_df(st.session_state.get("df", load_df_for(slug)))
    meta = load_meta_for(slug)
    label = st.selectbox("Sélectionnez un contrôle",
                         df.apply(lambda r: f"{r['Domain']} — {r['ID']} — {r['Item']}", axis=1),
                         key=f"evi_select_row_{slug}")
    if not label: return
    idx = df.index[df.apply(lambda x: f"{x['Domain']} — {x['ID']} — {x['Item']}", axis=1) == label][0]
    r = df.iloc[idx]; qid, item = r["ID"], r["Item"]
    st.write(f"**Contrôle :** {qid} — {item}")

    files = st.file_uploader("Ajouter des preuves", accept_multiple_files=True, key=f"evi_uploader_{slug}")
    cols = st.columns(3)
    if cols[0].button("Téléverser", key=f"evi_upload_btn_{slug}"):
        n = persist_uploads(meta.get("audit_id","AUDIT-001"), qid, item, files)
        st.success(f"{n} fichier(s) ajouté(s)."); st.rerun()

    st.write("**Pièces existantes :**")
    for p in list_evidence(meta.get("audit_id","AUDIT-001"), qid, item):
        c = st.columns([6,2,2])
        c[0].write(p.name)
        c[1].download_button("Télécharger", data=open(p, "rb").read(), file_name=p.name, key=f"evi_dl_{p.name}_{slug}")
        if c[2].button("Supprimer", key=f"evi_del_{p.name}_{slug}"):
            delete_evidence(p); st.rerun()

    if cols[1].button("📦 Exporter toutes les preuves (ZIP)", key=f"evi_zip_btn_{slug}"):
        z = export_evidence_zip(meta.get("audit_id","AUDIT-001"), df)
        st.download_button("Télécharger le ZIP", data=z,
                           file_name=f"evidence_{_slug(meta.get('audit_id','AUDIT-001'))}.zip",
                           key=f"evi_zip_dl_{slug}")

def _radar_by_domain(df: pd.DataFrame, out_path: Path) -> Path:
    d = df.copy(); d["Risk"] = d["Level"].map(_risk_score)
    by_dom = d.groupby("Domain")["Risk"].mean().sort_index() if len(d)>0 else pd.Series(dtype=float)
    labels = list(by_dom.index)
    values = (100 - (by_dom / 100.0) * 100).clip(0,100).values if len(by_dom)>0 else np.array([])

    fig = plt.figure()
    if len(labels)==0:
        plt.text(0.5,0.5,"Aucune donnée", ha="center", va="center")
    else:
        angles = np.linspace(0, 2*np.pi, len(labels), endpoint=False).tolist()
        values = np.concatenate((values, [values[0]])); angles += [angles[0]]
        ax = fig.add_subplot(111, polar=True)
        ax.plot(angles, values); ax.fill(angles, values, alpha=0.1)
        ax.set_thetagrids(np.degrees(angles[:-1]), labels)
        ax.set_title("Radar par domaine (plus haut = meilleur)"); ax.set_ylim(0,100)
    fig.savefig(out_path, bbox_inches="tight"); plt.close(fig); return out_path

def _heatmap_risques(df: pd.DataFrame, out_path: Path) -> Path:
    d = df.copy()
    d["RiskScore"] = d["Level"].map(_risk_from_level).map(_risk_score)
    cols_order = ["Élevé","Moyen","Conforme","Non applicable"]
    d["RiskName"] = d["Level"].map(_risk_from_level)
    pivot = d.pivot_table(index="Domain", columns="RiskName", values="RiskScore", aggfunc="sum", fill_value=0)
    for c in cols_order:
        if c not in pivot.columns: pivot[c]=0
    pivot = pivot[cols_order].sort_index()

    fig = plt.figure(); ax = fig.add_subplot(111); mat = ax.imshow(pivot.values, aspect="auto")
    ax.set_xticks(np.arange(len(pivot.columns))); ax.set_xticklabels(pivot.columns, rotation=45, ha="right")
    ax.set_yticks(np.arange(len(pivot.index))); ax.set_yticklabels(pivot.index)
    ax.set_title("Cartographie des risques (somme des scores)")
    for i in range(pivot.shape[0]):
        for j in range(pivot.shape[1]):
            ax.text(j, i, str(int(pivot.values[i, j])), ha="center", va="center")
    fig.colorbar(mat, ax=ax, shrink=0.7); fig.tight_layout(); fig.savefig(out_path, bbox_inches="tight"); plt.close(fig); return out_path

def _tab_synthese(slug: str):
    st.subheader("📊 Synthèse")
    df = _ensure_df(st.session_state.get("df", load_df_for(slug)))
    counts = df["Level"].value_counts()
    c1,c2,c3,c4 = st.columns(4)
    c1.metric("Conforme", int(counts.get("Conforme",0)))
    c2.metric("Partiellement conforme", int(counts.get("Partiellement conforme",0)))
    c3.metric("Non conforme", int(counts.get("Non conforme",0)))
    c4.metric("Non applicable", int(counts.get("Non applicable",0)))

    # Plan (édition rapide + autosave optionnel)
    st.sidebar.markdown("---")
    autosave_plan = _sidebar_toggle("Sauvegarde auto (Plan d'actions)", f"autosave_plan_{slug}", True)

    plan = st.session_state.get("plan_actions", load_plan_for(slug, df))
    st.session_state["plan_actions"] = plan

    if not plan.empty:
        a_faire = int((plan["Statut"]=="À faire").sum())
        en_cours = int((plan["Statut"]=="En cours").sum())
        clos = int((plan["Statut"]=="Clos").sum())
        effort_total = int(pd.to_numeric(plan.get("Effort (JH)"), errors="coerce").fillna(0).sum())
        budget_total = int(pd.to_numeric(plan.get("Budget estimé (€)"), errors="coerce").fillna(0).sum())
        avancement_moy = float(pd.to_numeric(plan.get("Avancement (%)"), errors="coerce").fillna(0).mean())

        s1,s2,s3,s4,s5 = st.columns(5)
        s1.metric("Actions (À faire)", a_faire)
        s2.metric("En cours", en_cours)
        s3.metric("Clôturées", clos)
        s4.metric("Effort total (JH)", effort_total)
        s5.metric("Budget total (€)", f"{budget_total:,}".replace(",", " "))

        try: st.progress(int(round(avancement_moy)), text=f"Avancement moyen du plan : {avancement_moy:.0f}%")
        except TypeError: st.progress(int(round(avancement_moy)))

    paths = _project_paths(slug)
    _heatmap_risques(df, paths["heatmap"]); _radar_by_domain(df, paths["radar"])

    cA, cB = st.columns(2)
    with cA: st.image(str(paths["heatmap"]), caption="Cartographie des risques", use_column_width=True)
    with cB: st.image(str(paths["radar"]), caption="Radar par domaine", use_column_width=True)

def _tab_exports(slug: str):
    st.subheader("📦 Exports & livrables")
    df = _ensure_df(st.session_state.get("df", load_df_for(slug)))
    plan = st.session_state.get("plan_actions", load_plan_for(slug, df))

    edited_plan = st.data_editor(
        plan,
        use_container_width=True,
        key=f"exp_plan_editor_{slug}",
        column_config={
            "Statut": st.column_config.SelectboxColumn(options=["À faire","En cours","Clos"]),
            "Effort (JH)": st.column_config.NumberColumn(min_value=0, step=1),
            "Budget estimé (€)": st.column_config.NumberColumn(min_value=0, step=100),
            "Avancement (%)": st.column_config.NumberColumn(min_value=0, max_value=100, step=5),
            "Échéance": st.column_config.TextColumn(help="YYYY-MM-DD"),
        },
        num_rows="dynamic",
    )
    if st.session_state.get(f"autosave_plan_{slug}", True):
        new_sig = _sig_csv(edited_plan)
        last_key = f"plan_last_sig_{slug}"
        if st.session_state.get(last_key) != new_sig:
            st.session_state["plan_actions"] = edited_plan
            save_plan_for(slug, edited_plan)
            st.session_state[last_key] = new_sig
            try: st.toast("Plan d'actions sauvegardé automatiquement.", icon="💾")
            except Exception: st.info("Plan d'actions sauvegardé (auto).")
    else:
        st.session_state["plan_actions"] = edited_plan

    cols = st.columns(4)
    if cols[0].button("💾 Sauver le plan (projet)", key=f"exp_plan_save_btn_{slug}"):
        save_plan_for(slug, st.session_state["plan_actions"]); st.success("Plan enregistré.")

    if cols[1].button("⬇️ Exporter dataset (CSV)", key=f"exp_data_btn_{slug}"):
        st.download_button("Télécharger les données (CSV)",
                           data=df.to_csv(index=False).encode("utf-8"),
                           file_name=f"dataset_{slug}.csv",
                           key=f"exp_data_dl_{slug}")

    meta = load_meta_for(slug)
    client  = cols[2].text_input("Client", value=meta.get("client","Client X"), key=f"exp_client_{slug}")
    auditor = cols[3].text_input("Auditeur", value=meta.get("auditor",""), key=f"exp_auditor_{slug}")
    rdate   = st.date_input("Date du rapport", value=pd.to_datetime(meta.get("date", datetime.utcnow().date())).date(), key=f"exp_date_{slug}")
    meta["client"]=client; meta["auditor"]=auditor; meta["date"]=str(rdate); save_meta_for(slug, meta)

    paths = _project_paths(slug)
    _heatmap_risques(df, paths["heatmap"]); _radar_by_domain(df, paths["radar"])

    if st.button("📄 Générer le rapport ISACA (DOCX)", type="primary", key=f"exp_docx_btn_{slug}"):
        try:
            generate_isaca_docx(meta.get("audit_id","AUDIT-001"), df, meta, st.session_state["plan_actions"],
                                {"heatmap": paths["heatmap"], "radar": paths["radar"]},
                                paths["docx"], logo_path=(paths["logo"] if paths["logo"].exists() else None))
            st.success("Rapport généré.")
            st.download_button("Télécharger le rapport ISACA", data=open(paths["docx"],"rb").read(),
                               file_name=paths["docx"].name,
                               mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                               key=f"exp_docx_dl_{slug}")
        except Exception as e:
            st.error(f"Erreur génération DOCX: {e}")

    # PowerPoint premium
    brand_hex = st.session_state.get("settings", DEFAULT_SETTINGS).get("brand_color_hex","2F5597")
    if Presentation is None:
        st.button("🖼️ Générer la présentation PowerPoint (PPTX)", disabled=True, help="Installer: pip install python-pptx")
        st.caption("Astuce : `pip install python-pptx` puis relance l’app pour activer l’export PPTX.")
    else:
        if st.button("🖼️ Générer la présentation PowerPoint (PPTX)", key=f"exp_pptx_btn_{slug}"):
            try:
                p = generate_pptx_summary(meta.get("audit_id","AUDIT-001"), meta, df, st.session_state["plan_actions"],
                                          {"heatmap": paths["heatmap"], "radar": paths["radar"]}, paths["pptx"], brand_hex)
                st.success("Présentation générée.")
                st.download_button("Télécharger la présentation", data=open(paths["pptx"],"rb").read(),
                                   file_name=paths["pptx"].name, key=f"exp_pptx_dl_{slug}")
            except Exception as e:
                st.error(f"Erreur génération PPTX: {e}")

def _tab_admin(slug: str):
    st.subheader("🛡️ Admin — Normes & Paramètres & Utilisateurs")
    tab_load, tab_edit, tab_catalog, tab_settings, tab_users = st.tabs(["Charger", "Éditeur", "Catalogue", "Paramétrage", "Utilisateurs"])

    # Charger norme
    with tab_load:
        st.markdown("#### Charger une norme depuis un fichier")
        nname = st.text_input("Nom (alias) de la norme", key=f"admin_norm_alias_{slug}")
        up = st.file_uploader("Fichier norme (CSV/XLSX/JSON)", key=f"admin_norm_upload_{slug}")
        if st.button("📤 Enregistrer la norme", key=f"admin_norm_save_{slug}"):
            if not nname or not up:
                st.error("Nom et fichier requis.")
            else:
                try:
                    name = (up.name or "").lower()
                    if name.endswith(".csv"):
                        df = pd.read_csv(up, sep=None, engine="python")
                    elif name.endswith((".xlsx",".xls")):
                        df = pd.read_excel(up)
                    elif name.endswith(".json"):
                        df = pd.read_json(up)
                    else:
                        raise ValueError("Format non supporté (CSV/XLSX/JSON).")
                    p = _save_norm_df(nname, df)
                    st.success(f"Norme enregistrée: {p.name}")
                    # appliquer au projet
                    st.session_state["df"] = _ensure_df(df)
                    meta = load_meta_for(slug); meta["norm"] = _slug(nname); save_meta_for(slug, meta)
                    save_df_for(slug, st.session_state["df"])
                    _set_active_norm(nname)
                    st.success(f"Norme '{nname}' appliquée et persistée.")
                    st.rerun()
                except Exception as e:
                    st.error(f"Erreur import: {e}")

        st.markdown("#### Normes disponibles")
        norms = _list_norms()
        if norms:
            sel = st.radio("Sélection active (globale)", norms,
                           index=norms.index(_get_active_norm()) if _get_active_norm() in norms else 0,
                           key=f"admin_norm_active_{slug}")
            cc = st.columns(4)
            if cc[0].button("✅ Activer globalement", key=f"admin_norm_activate_{slug}"):
                try:
                    _set_active_norm(sel)
                    st.session_state["df"] = _load_norm_df(sel)
                    meta = load_meta_for(slug); meta["norm"] = sel; save_meta_for(slug, meta)
                    save_df_for(slug, st.session_state["df"])
                    st.success(f"Norme active et appliquée au projet : {sel}")
                    st.rerun()
                except Exception as e:
                    st.error(f"Activation impossible: {e}")
            if cc[1].button("🗑️ Supprimer la norme", key=f"admin_norm_delete_{slug}"):
                try:
                    (NORMS_DIR / f"{_slug(sel)}.csv").unlink(missing_ok=True)
                    if _get_active_norm() == sel: ACTIVE_NORM.unlink(missing_ok=True)
                    st.success("Norme supprimée."); st.rerun()
                except Exception as e:
                    st.error(f"Suppression impossible: {e}")
            if cc[2].button("👁️ Charger dans l'éditeur", key=f"admin_norm_open_editor_{slug}"):
                st.session_state[NORM_EDIT_ALIAS_KEY] = sel
                try:
                    st.session_state[NORM_EDIT_DATA_KEY] = _ensure_df(_load_norm_df(sel))
                    st.success(f"Norme '{sel}' chargée dans l'éditeur."); st.rerun()
                except Exception as e:
                    st.error(f"Erreur de chargement: {e}")
            if cc[3].button("📚 Appliquer au projet courant", key=f"admin_apply_to_proj_{slug}"):
                try:
                    st.session_state["df"] = _load_norm_df(sel)
                    meta = load_meta_for(slug); meta["norm"] = sel; save_meta_for(slug, meta)
                    save_df_for(slug, st.session_state["df"])
                    st.success(f"Norme '{sel}' appliquée et persistée."); st.rerun()
                except Exception as e:
                    st.error(f"Impossible d'appliquer : {e}")
        else:
            st.info("Aucune norme enregistrée.")

    # Éditeur norme
    with tab_edit:
        st.markdown("#### Éditeur de norme")
        alias = st.text_input("Alias à éditer (nouveau ou existant)",
                              value=st.session_state.get(NORM_EDIT_ALIAS_KEY, ""),
                              key=f"norm_edit_alias_{slug}")
        bnew, bload, bsave, bact = st.columns(4)
        if bnew.button("🆕 Nouveau (vide)", key=f"norm_edit_new_{slug}"):
            st.session_state[NORM_EDIT_ALIAS_KEY] = alias or "nouvelle_norme"
            st.session_state[NORM_EDIT_DATA_KEY] = _ensure_df(pd.DataFrame(columns=REQUIRED_COLS))
            st.rerun()
        if bload.button("📥 Charger existant", key=f"norm_edit_load_{slug}"):
            try:
                st.session_state[NORM_EDIT_ALIAS_KEY] = alias
                st.session_state[NORM_EDIT_DATA_KEY]  = _load_norm_df(alias)
                st.rerun()
            except Exception as e:
                st.error(f"Chargement impossible: {e}")

        st.info("Colonnes requises : " + ", ".join(REQUIRED_COLS))
        base_df = _ensure_df(st.session_state.get(NORM_EDIT_DATA_KEY, pd.DataFrame(columns=REQUIRED_COLS)))

        edited = st.data_editor(
            base_df, use_container_width=True, num_rows="dynamic",
            key=NORM_EDIT_WIDGET_KEY,
            column_config={"Level": st.column_config.SelectboxColumn(options=EVAL_OPTIONS, required=True)},
        )

        if st.button("Appliquer les modifs (mémoire)", key=f"norm_edit_apply_{slug}"):
            st.session_state[NORM_EDIT_DATA_KEY] = edited; st.success("Modifs appliquées.")
        if bsave.button("💾 Sauvegarder (écrase)", key=f"norm_edit_save_{slug}"):
            try:
                to_save = _ensure_df(edited)
                st.session_state[NORM_EDIT_DATA_KEY] = to_save.copy()
                p = _save_norm_df(st.session_state.get(NORM_EDIT_ALIAS_KEY) or alias, to_save)
                st.success(f"Norme sauvegardée: {p.name}")
            except Exception as e:
                st.error(f"Sauvegarde impossible: {e}")
        if bact.button("✅ Activer globalement", key=f"norm_edit_activate_{slug}"):
            try:
                _set_active_norm(st.session_state.get(NORM_EDIT_ALIAS_KEY) or alias)
                st.success(f"Norme active: {st.session_state.get(NORM_EDIT_ALIAS_KEY) or alias}")
            except Exception as e:
                st.error(f"Activation impossible: {e}")

    # Catalogue
    with tab_catalog:
        st.markdown("#### Présets (squelettes rapides)")
        c1, c2, c3 = st.columns(3)
        alias_iso  = c1.text_input("Alias ISO 27001", value="iso-27001", key=f"cat_iso_alias_{slug}")
        alias_nis2 = c2.text_input("Alias NIS2",     value="nis2",      key=f"cat_nis2_alias_{slug}")
        alias_dora = c3.text_input("Alias DORA",     value="dora",      key=f"cat_dora_alias_{slug}")
        if c1.button("Créer squelette ISO 27001", key=f"cat_iso_btn_{slug}"):
            st.success(f"Créé: {_save_norm_df(alias_iso,  _ensure_df(pd.DataFrame([{'Domain':'Organisation','ID':'ORG-01','Item':'Gouvernance','Contrôle':'Rôles & responsabilités','Level':'Partiellement conforme','Comment':''}]))).name}")
        if c2.button("Créer squelette NIS2", key=f"cat_nis2_btn_{slug}"):
            st.success(f"Créé: {_save_norm_df(alias_nis2, _ensure_df(pd.DataFrame([{'Domain':'Gestion du risque','ID':'RIS-01','Item':'Méthodo risque','Contrôle':'Méthodologie formalisée','Level':'Partiellement conforme','Comment':''}]))).name}")
        if c3.button("Créer squelette DORA", key=f"cat_dora_btn_{slug}"):
            st.success(f"Créé: {_save_norm_df(alias_dora, _ensure_df(pd.DataFrame([{'Domain':'ICT Risk','ID':'ICT-01','Item':'Gouv ICT','Contrôle':'Cadre de gestion des risques','Level':'Partiellement conforme','Comment':''}]))).name}")

    # Paramétrage
    with tab_settings:
        st.markdown("#### Paramétrage global")
        settings = st.session_state.get("settings", load_settings())

        exp1 = st.expander("Niveaux & Risques", expanded=True)
        with exp1:
            c1, c2 = st.columns(2)
            lvl_json = c1.text_area("Normalisation (JSON)", value=json.dumps(settings.get("level_normalization", {}), ensure_ascii=False, indent=2), key=f"set_norm_json_{slug}", height=220)
            # Scores/échéances par niveau de risque interne
            df_score = pd.DataFrame(list(settings.get("risk_score", {}).items()), columns=["Niveau de risque","Score"]).sort_values(by="Score", ascending=False)
            df_deadl = pd.DataFrame(list(settings.get("deadlines_days", {}).items()), columns=["Niveau de risque","Jours"]).sort_values(by="Jours", ascending=True)
            df_score = st.data_editor(df_score, key=f"set_scores_{slug}", num_rows="dynamic", use_container_width=True)
            df_deadl = st.data_editor(df_deadl, key=f"set_deadlines_{slug}", num_rows="dynamic", use_container_width=True)

        exp2 = st.expander("Chiffrage & Ressources", expanded=False)
        with exp2:
            c3, c4 = st.columns(2)
            df_eff = pd.DataFrame(list(settings.get("effort_days", {}).items()), columns=["Niveau de risque","Effort (JH)"]).sort_values(by="Effort (JH)", ascending=False)
            df_bud = pd.DataFrame(list(settings.get("budget_eur", {}).items()), columns=["Niveau de risque","Budget (€)"]).sort_values(by="Budget (€)", ascending=False)
            df_eff = st.data_editor(df_eff, key=f"set_effort_{slug}", num_rows="dynamic", use_container_width=True)
            df_bud = st.data_editor(df_bud, key=f"set_budget_{slug}", num_rows="dynamic", use_container_width=True)
            st.write("Profils de ressources (motif regex → humains/techniques)")
            prof_df = pd.DataFrame(settings.get("resource_profiles", []))
            prof_df = st.data_editor(prof_df, key=f"set_profiles_{slug}", num_rows="dynamic", use_container_width=True)

        exp3 = st.expander("Recommandations & KPI", expanded=False)
        with exp3:
            recs_json = st.text_area("Recommandations (JSON)", value=json.dumps(settings.get("recs", {}), ensure_ascii=False, indent=2), key=f"set_recs_json_{slug}", height=300)

        csave, creset, cexp, cimp = st.columns(4)
        if csave.button("💾 Enregistrer paramètres", key=f"set_save_btn_{slug}"):
            try:
                settings["level_normalization"] = json.loads(lvl_json)
                settings["risk_score"]     = {str(k): int(v) for k, v in pd.DataFrame(df_score).dropna().values}
                settings["deadlines_days"] = {str(k): int(v) for k, v in pd.DataFrame(df_deadl).dropna().values}
                settings["effort_days"]    = {str(k): int(v) for k, v in pd.DataFrame(df_eff).dropna().values}
                settings["budget_eur"]     = {str(k): int(v) for k, v in pd.DataFrame(df_bud).dropna().values}
                settings["resource_profiles"] = [dict(r) for _, r in pd.DataFrame(prof_df).fillna("").iterrows()]
                settings["recs"] = json.loads(recs_json)
                save_settings(settings); st.session_state["settings"] = settings
                st.success("Paramètres enregistrés ✅")
            except Exception as e:
                st.error(f"Impossible d'enregistrer : {e}")

        if creset.button("↩️ Valeurs par défaut", type="secondary", key=f"set_reset_btn_{slug}"):
            save_settings(DEFAULT_SETTINGS); st.session_state["settings"] = load_settings()
            st.success("Paramètres réinitialisés.")

        if cexp.button("⬇️ Export paramètres (Excel)", key=f"set_export_btn_{slug}"):
            try:
                data = export_settings_to_excel(st.session_state.get("settings", settings))
                st.download_button("Télécharger settings.xlsx", data=data, file_name="settings.xlsx", key=f"set_export_dl_{slug}")
            except Exception as e:
                st.error(f"Export Excel impossible (openpyxl manquant ?): {e}")

        imp_file = cimp.file_uploader("Importer paramètres (Excel)", type=["xlsx"], key=f"set_import_file_{slug}")
        if imp_file is not None and cimp.button("Importer", key=f"set_import_btn_{slug}"):
            try:
                new_s = import_settings_from_excel(imp_file)
                save_settings(new_s); st.session_state["settings"] = new_s
                st.success("Paramètres importés ✅"); st.rerun()
            except Exception as e:
                st.error(f"Import impossible : {e}")

    # Utilisateurs
    with tab_users:
        st.markdown("#### Utilisateurs")
        u1,u2 = st.columns(2)
        with u1:
            n_email = st.text_input("Email (nouvel utilisateur)", key=f"admin_user_email_new_{slug}")
            n_name  = st.text_input("Nom complet", key=f"admin_user_fullname_new_{slug}")
            n_pwd   = st.text_input("Mot de passe", type="password", key=f"admin_user_pwd_new_{slug}")
            if st.button("Créer l'utilisateur", key=f"admin_user_create_{slug}"):
                try:
                    res = AuthAdapter.create_user(n_email, n_pwd, n_name or n_email, role="user", is_active=True)
                    if isinstance(res, tuple) and res and res[0] is False: st.error(res[1] or "Création refusée")
                    else: st.success("Utilisateur créé.")
                except Exception as e:
                    st.error(f"Erreur création utilisateur: {e}")
        with u2:
            c_email = st.text_input("Email (modifier le mot de passe)", key=f"admin_user_email_change_{slug}")
            c_pwd   = st.text_input("Nouveau mot de passe", type="password", key=f"admin_user_pwd_change_{slug}")
            if st.button("Changer le mot de passe", key=f"admin_user_change_btn_{slug}"):
                try:
                    if hasattr(auth, "set_password"):
                        auth.set_password(c_email, c_pwd); st.success("Mot de passe modifié.")
                    else:
                        st.error("set_password indisponible dans auth.py")
                except Exception as e:
                    st.error(f"Erreur: {e}")

# ------------------------ Main ------------------------
def main():
    st.set_page_config(page_title="CyberPivot", page_icon="🛡️", layout="wide")
    st.title("CyberPivot — Missions Audit (ISACA+)")
    st.caption(f"Build: {APP_BUILD} • {datetime.utcnow().isoformat()}Z")
    try: auth.init_auth_db()
    except Exception: pass

    if "settings" not in st.session_state:
        st.session_state["settings"] = load_settings()

    # Projet courant
    slug = st.session_state.get("project_slug") or "audit-001"
    slug = _slug(slug)
    st.session_state["project_slug"] = slug
    _touch_recent(slug)

    # Charger meta + df + plan
    meta = load_meta_for(slug)
    st.session_state["df"] = _ensure_df(load_df_for(slug, meta))
    st.session_state["plan_actions"] = load_plan_for(slug, st.session_state["df"])

    # Auth
    user = st.session_state.get("_auth") if isinstance(st.session_state.get("_auth"), dict) else None
    if not user or not user.get("status"):
        st.header("🔐 Connexion")
        tab_co, tab_new = st.tabs(["Se connecter", "Créer un compte"])
        with tab_co:
            email = st.text_input("Email", key="login_email")
            pwd   = st.text_input("Mot de passe", type="password", key="login_pwd")
            if st.button("Connexion", type="primary", key="login_btn"):
                u = AuthAdapter.login_get_user(email, pwd)
                if not u: st.error("Identifiants invalides."); st.stop()
                if not bool(u.get("is_active", u.get("active", 1))): st.error("Compte inactif."); st.stop()
                st.session_state["_auth"] = {"email": u.get("email", email), "full_name": u.get("full_name") or u.get("name") or email, "role": u.get("role","user"), "status": True}
                st.rerun()
        with tab_new:
            n_email = st.text_input("Email (nouveau)", key="signup_email")
            n_name  = st.text_input("Nom complet", key="signup_fullname")
            n_pwd1  = st.text_input("Mot de passe", type="password", key="signup_pwd1")
            n_pwd2  = st.text_input("Confirmer", type="password", key="signup_pwd2")
            if st.button("Créer le compte", key="signup_btn"):
                if not n_email or not n_pwd1: st.error("Email et mot de passe requis."); st.stop()
                if n_pwd1 != n_pwd2: st.error("Les mots de passe ne correspondent pas."); st.stop()
                res = AuthAdapter.create_user(n_email, n_pwd1, n_name or n_email, role="user", is_active=True)
                if isinstance(res, tuple) and res and res[0] is False: st.error(res[1] or "Création refusée")
                else: st.success("Compte créé. Connectez-vous.")
        st.stop()

    # Top bar
    _topbar(user, slug)

    tabs = ["Infos client", "Audit", "Preuves", "Synthèse", "Exports"]
    if user.get("role","user") == "admin": tabs.append("Admin")
    t_objs = st.tabs(tabs)

    with t_objs[0]: _tab_infos_client(slug)
    with t_objs[1]: _tab_audit(slug)
    with t_objs[2]: _tab_preuves(slug)
    with t_objs[3]: _tab_synthese(slug)
    with t_objs[4]: _tab_exports(slug)
    if user.get("role","user") == "admin":
        with t_objs[5]: _tab_admin(slug)

if __name__ == "__main__":
    main()

